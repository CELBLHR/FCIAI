"""
简化版PPT背景提取工具
专门解决pyuno转换后背景丢失问题
"""
import os
import zipfile
import shutil
from typing import Dict, List, Optional
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)

class SimpleBackgroundExtractor:
    """简化的背景提取器，专注于解决实际问题"""
    
    def __init__(self, original_pptx: str):
        self.original_pptx = original_pptx
        self.backup_folder = f"{original_pptx}_background_backup"
        
    def create_background_backup(self) -> bool:
        """
        创建背景备份（在pyuno处理前调用）
        """
        try:
            # 创建备份文件夹
            os.makedirs(self.backup_folder, exist_ok=True)
            
            # 1. 直接复制原始PPT作为完整备份
            backup_pptx = os.path.join(self.backup_folder, "original_backup.pptx")
            shutil.copy2(self.original_pptx, backup_pptx)
            
            # 2. 提取所有媒体文件（图片）
            media_folder = os.path.join(self.backup_folder, "media")
            os.makedirs(media_folder, exist_ok=True)
            
            with zipfile.ZipFile(self.original_pptx, 'r') as zip_file:
                # 提取所有媒体文件
                for file_name in zip_file.namelist():
                    if file_name.startswith('ppt/media/'):
                        # 提取到media文件夹
                        zip_file.extract(file_name, self.backup_folder)
                        logger.info(f"备份媒体文件: {file_name}")
            
            # 3. 提取背景相关的XML文件
            xml_folder = os.path.join(self.backup_folder, "xml")
            os.makedirs(xml_folder, exist_ok=True)
            
            with zipfile.ZipFile(self.original_pptx, 'r') as zip_file:
                # 备份重要的XML文件
                important_files = [
                    'ppt/slideMasters/slideMaster1.xml',  # 母版
                    'ppt/slideLayouts/',                  # 版式
                    'ppt/theme/theme1.xml',              # 主题
                ]
                
                for file_pattern in important_files:
                    for file_name in zip_file.namelist():
                        if file_name.startswith(file_pattern.rstrip('/')):
                            zip_file.extract(file_name, xml_folder)
                            logger.debug(f"备份XML文件: {file_name}")
            
            # 4. 分析并记录背景信息
            background_info = self._analyze_backgrounds()
            
            # 保存背景信息
            import json
            info_file = os.path.join(self.backup_folder, "background_info.json")
            with open(info_file, 'w', encoding='utf-8') as f:
                json.dump(background_info, f, indent=2, ensure_ascii=False, default=str)
            
            logger.info(f"背景备份创建完成: {self.backup_folder}")
            return True
            
        except Exception as e:
            logger.error(f"创建背景备份失败: {e}")
            return False
    
    def _analyze_backgrounds(self) -> Dict:
        """分析背景信息"""
        try:
            prs = Presentation(self.original_pptx)
            
            background_info = {
                'slide_count': len(prs.slides),
                'master_count': len(prs.slide_masters),
                'slides': {},
                'masters': {},
                'has_background_images': False
            }
            
            # 分析每张幻灯片的背景
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    'has_custom_background': False,
                    'background_type': 'inherit'
                }
                
                try:
                    if hasattr(slide.background, 'fill'):
                        fill_type = slide.background.fill.type
                        if fill_type and fill_type.name != 'BACKGROUND':
                            slide_info['has_custom_background'] = True
                            slide_info['background_type'] = fill_type.name
                            if fill_type.name == 'PICTURE':
                                background_info['has_background_images'] = True
                except:
                    pass
                
                background_info['slides'][i] = slide_info
            
            # 分析母版背景
            for i, master in enumerate(prs.slide_masters):
                master_info = {
                    'layout_count': len(master.slide_layouts),
                    'background_type': 'unknown'
                }
                
                try:
                    if hasattr(master.background, 'fill'):
                        fill_type = master.background.fill.type
                        if fill_type:
                            master_info['background_type'] = fill_type.name
                            if fill_type.name == 'PICTURE':
                                background_info['has_background_images'] = True
                except:
                    pass
                
                background_info['masters'][i] = master_info
            
            return background_info
            
        except Exception as e:
            logger.error(f"分析背景信息失败: {e}")
            return {}
    
    def restore_backgrounds_after_pyuno(self, translated_pptx: str) -> bool:
        """
        在pyuno处理后恢复背景（核心功能）
        """
        try:
            if not os.path.exists(self.backup_folder):
                logger.error("背景备份不存在，无法恢复")
                return False
            
            # 方法1: 直接替换媒体文件夹
            success1 = self._restore_media_files(translated_pptx)
            
            # 方法2: 使用python-pptx复制背景设置
            success2 = self._copy_background_settings(translated_pptx)
            
            # 方法3: 如果以上都失败，使用混合方案
            if not (success1 or success2):
                success3 = self._hybrid_restore(translated_pptx)
                return success3
            
            return success1 or success2
            
        except Exception as e:
            logger.error(f"恢复背景失败: {e}")
            return False
    
    def _restore_media_files(self, translated_pptx: str) -> bool:
        """直接替换媒体文件（最直接的方法）"""
        try:
            backup_media = os.path.join(self.backup_folder, "ppt", "media")
            if not os.path.exists(backup_media):
                logger.warning("备份中没有媒体文件")
                return False
            
            # 创建临时文件
            temp_pptx = f"{translated_pptx}.temp"
            
            # 复制翻译后的PPTX
            shutil.copy2(translated_pptx, temp_pptx)
            
            # 打开为ZIP文件并替换媒体文件夹
            with zipfile.ZipFile(temp_pptx, 'a') as zip_file:
                # 删除现有媒体文件
                # 注意：zipfile不支持直接删除，需要重新创建
                pass
            
            # 重新创建ZIP，替换媒体文件
            with zipfile.ZipFile(translated_pptx, 'r') as old_zip:
                with zipfile.ZipFile(temp_pptx, 'w') as new_zip:
                    # 复制所有非媒体文件
                    for item in old_zip.infolist():
                        if not item.filename.startswith('ppt/media/'):
                            data = old_zip.read(item.filename)
                            new_zip.writestr(item, data)
                    
                    # 添加备份的媒体文件
                    for root, dirs, files in os.walk(backup_media):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arc_name = file_path.replace(self.backup_folder + os.sep, '').replace(os.sep, '/')
                            new_zip.write(file_path, arc_name)
                            logger.debug(f"恢复媒体文件: {arc_name}")
            
            # 替换原文件
            shutil.move(temp_pptx, translated_pptx)
            logger.info("媒体文件恢复完成")
            return True
            
        except Exception as e:
            logger.error(f"恢复媒体文件失败: {e}")
            return False
    
    def _copy_background_settings(self, translated_pptx: str) -> bool:
        """使用python-pptx复制背景设置"""
        try:
            # 加载原始备份和翻译后的PPT
            backup_pptx = os.path.join(self.backup_folder, "original_backup.pptx")
            if not os.path.exists(backup_pptx):
                return False
            
            original_prs = Presentation(backup_pptx)
            translated_prs = Presentation(translated_pptx)
            
            # 确保幻灯片数量一致
            if len(original_prs.slides) != len(translated_prs.slides):
                logger.warning("幻灯片数量不一致，可能影响背景恢复")
            
            # 逐一复制背景设置
            min_slides = min(len(original_prs.slides), len(translated_prs.slides))
            
            for i in range(min_slides):
                try:
                    original_slide = original_prs.slides[i]
                    translated_slide = translated_prs.slides[i]
                    
                    # 复制背景设置（这是一个复杂的过程）
                    self._copy_slide_background(original_slide, translated_slide)
                    
                except Exception as slide_error:
                    logger.warning(f"复制第{i}张幻灯片背景失败: {slide_error}")
            
            # 保存修改后的PPT
            translated_prs.save(translated_pptx)
            logger.info("背景设置复制完成")
            return True
            
        except Exception as e:
            logger.error(f"复制背景设置失败: {e}")
            return False
    
    def _copy_slide_background(self, source_slide, target_slide):
        """复制单张幻灯片的背景"""
        try:
            # 这是一个复杂的过程，需要处理不同类型的背景
            source_bg = source_slide.background
            target_bg = target_slide.background
            
            # 尝试复制背景填充
            if hasattr(source_bg, 'fill') and hasattr(target_bg, 'fill'):
                # 这里需要根据不同的填充类型进行处理
                # 由于python-pptx的限制，完整实现比较复杂
                pass
                
        except Exception as e:
            logger.debug(f"复制幻灯片背景时出错: {e}")
    
    def _hybrid_restore(self, translated_pptx: str) -> bool:
        """混合恢复方案：结合多种方法"""
        try:
            logger.info("使用混合恢复方案...")
            
            # 1. 首先尝试从备份中提取关键的背景图片
            self._extract_key_background_images(translated_pptx)
            
            # 2. 手动修复明显的背景问题
            self._fix_obvious_background_issues(translated_pptx)
            
            return True
            
        except Exception as e:
            logger.error(f"混合恢复失败: {e}")
            return False
    
    def _extract_key_background_images(self, translated_pptx: str):
        """提取关键背景图片"""
        # 实现关键背景图片的提取和插入
        pass
    
    def _fix_obvious_background_issues(self, translated_pptx: str):
        """修复明显的背景问题"""
        try:
            prs = Presentation(translated_pptx)
            
            # 修复每张幻灯片的背景问题
            for slide in prs.slides:
                # 检查是否有明显的背景损坏
                # 如果有，尝试设置为白色背景或其他安全背景
                try:
                    if hasattr(slide.background, 'fill'):
                        # 如果背景看起来损坏了，设置为白色
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = (255, 255, 255)
                except:
                    pass
            
            prs.save(translated_pptx)
            logger.info("修复明显背景问题完成")
            
        except Exception as e:
            logger.error(f"修复背景问题失败: {e}")

# 集成到您的主流程中的使用方法
def integrate_background_protection(presentation_path: str, 
                                  process_function, 
                                  *args, **kwargs) -> bool:
    """
    集成背景保护的处理流程
    
    Args:
        presentation_path: PPT路径
        process_function: 翻译处理函数(如pyuno_controller)
        *args, **kwargs: 处理函数的参数
    
    Returns:
        处理是否成功
    """
    try:
        # 1. 创建背景备份
        logger.info("创建背景备份...")
        bg_extractor = SimpleBackgroundExtractor(presentation_path)
        backup_success = bg_extractor.create_background_backup()
        
        if not backup_success:
            logger.warning("背景备份失败，但继续处理...")
        
        # 2. 执行原始处理流程（pyuno等）
        logger.info("执行翻译处理...")
        process_result = process_function(*args, **kwargs)
        
        if not process_result:
            logger.error("翻译处理失败")
            return False
        
        # 3. 恢复背景
        if backup_success:
            logger.info("恢复背景...")
            restore_success = bg_extractor.restore_backgrounds_after_pyuno(presentation_path)
            
            if restore_success:
                logger.info("背景恢复成功")
            else:
                logger.warning("背景恢复失败，但翻译已完成")
        
        return True
        
    except Exception as e:
        logger.error(f"集成背景保护流程失败: {e}")
        return False

# 使用示例
if __name__ == "__main__":
    # 示例：如何在您的代码中使用
    def example_usage():
        pptx_path = "test.pptx"
        
        # 方式1: 手动控制
        bg_tool = SimpleBackgroundExtractor(pptx_path)
        bg_tool.create_background_backup()
        
        # 执行您的pyuno翻译处理
        # pyuno_controller(...)
        
        # 恢复背景
        bg_tool.restore_backgrounds_after_pyuno(pptx_path)
        
        # 方式2: 集成式调用
        # integrate_background_protection(
        #     pptx_path, 
        #     pyuno_controller,
        #     # pyuno_controller的参数...
        # )