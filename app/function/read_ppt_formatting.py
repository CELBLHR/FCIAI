import os
import logging
import uno
from com.sun.star.beans import PropertyValue
from com.sun.star.drawing import FillStyle, LineStyle
from com.sun.star.text import TextContentAnchorType
from com.sun.star.awt import FontSlant, FontWeight
from typing import Dict, List, Any, Optional

# 导入 LibreOfficeUNOColorManager
from .libreoffice_uno_color import LibreOfficeUNOColorManager, UNO_AVAILABLE

logger = logging.getLogger(__name__)

def get_ppt_format_uno(ppt_path: str, slide_index: Optional[int] = None) -> Dict[str, Any]:
    """
    使用LibreOffice UNO接口获取PPT文件中指定页面的格式信息。

    :param ppt_path: PowerPoint文件的路径。
    :param slide_index: 目标幻灯片的索引（从0开始）。如果为None，则获取所有幻灯片的信息。
    :return: 一个包含幻灯片及其所有形状格式信息的字典。
    """
    if not UNO_AVAILABLE:
        logger.error("LibreOffice UNO接口不可用，请安装LibreOffice并配置Python UNO SDK。")
        return {"error": "LibreOffice UNO接口不可用"}

    if not os.path.exists(ppt_path):
        logger.error(f"PPT文件不存在: {ppt_path}")
        return {"error": f"PPT文件不存在: {ppt_path}"}

    manager = LibreOfficeUNOColorManager()
    format_info = {"slides": []}

    try:
        logger.info(f"开始使用UNO获取PPT格式信息: {os.path.basename(ppt_path)}")

        # 1. 启动LibreOffice服务
        if not manager.start_libreoffice_service():
            logger.error("启动LibreOffice UNO服务失败")
            return {"error": "启动LibreOffice UNO服务失败"}

        # 2. 打开PPT文件
        if not manager.open_presentation(ppt_path):
            logger.error("打开PPT文件失败")
            return {"error": "打开PPT文件失败"}

        # 获取文档的DrawPages（幻灯片集合）
        draw_pages = manager.document.getDrawPages()
        total_slides = draw_pages.getCount()

        slides_to_process = []
        if slide_index is not None:
            if 0 <= slide_index < total_slides:
                slides_to_process.append(slide_index)
            else:
                logger.warning(f"幻灯片索引 {slide_index} 超出范围。总共有 {total_slides} 页。")
                return {"error": f"幻灯片索引 {slide_index} 超出范围。"}
        else:
            slides_to_process = list(range(total_slides))

        for idx in slides_to_process:
            slide = draw_pages.getByIndex(idx)
            slide_info = {
                "slide_index": idx,
                "width_cm": slide.Width / 1000,  # 转换为厘米 (1/100 mm)
                "height_cm": slide.Height / 1000, # 转换为厘米 (1/100 mm)
                "shapes": []
            }

            # 遍历幻灯片中的所有形状
            for shape_idx in range(slide.getCount()):
                shape = slide.getByIndex(shape_idx)
                shape_data = _extract_shape_info(shape, idx, shape_idx)
                if shape_data:
                    slide_info["shapes"].append(shape_data)
            
            format_info["slides"].append(slide_info)

        logger.info("✅ UNO获取PPT格式信息完成。")
        return format_info

    except Exception as e:
        logger.error(f"UNO获取PPT格式信息过程中出错: {e}", exc_info=True)
        return {"error": f"UNO获取PPT格式信息过程中出错: {e}"}
    finally:
        manager.cleanup()

def _extract_shape_info(shape: Any, slide_idx: int, shape_idx: int) -> Optional[Dict[str, Any]]:
    """
    提取单个形状的格式信息。
    """
    shape_info = {
        "slide_index": slide_idx,
        "shape_index": shape_idx,
        "type": _get_shape_type_name(shape),
        "name": shape.Name,
        "left_cm": shape.Position.X / 1000,  # 转换为厘米
        "top_cm": shape.Position.Y / 1000,   # 转换为厘米
        "width_cm": shape.Size.Width / 1000, # 转换为厘米
        "height_cm": shape.Size.Height / 1000, # 转换为厘米
        "rotation_deg": shape.RotateAngle / 100, # 转换为度
    }

    # 提取文本框信息
    if hasattr(shape, 'Text'):
        shape_info["text_content"] = shape.Text
        shape_info["text_frame"] = _extract_text_frame_info(shape)

    # 提取填充信息
    if hasattr(shape, 'FillStyle'):
        shape_info["fill"] = _extract_fill_info(shape)

    # 提取线条信息
    if hasattr(shape, 'LineStyle'):
        shape_info["line"] = _extract_line_info(shape)
    
    # 提取图片信息 (如果形状是图片)
    if hasattr(shape, 'GraphicURL') and shape.GraphicURL:
        shape_info["graphic_url"] = shape.GraphicURL
        shape_info["graphic_filter"] = shape.GraphicFilter

    return shape_info

def _get_shape_type_name(shape: Any) -> str:
    """根据UNO形状服务名称获取更友好的形状类型名称。"""
    try:
        # UNO形状的服务名称通常以 "com.sun.star.drawing." 开头
        service_name = shape.getServiceName()
        if "TextShape" in service_name:
            return "TEXT_BOX"
        elif "GraphicObjectShape" in service_name:
            return "PICTURE"
        elif "RectangleShape" in service_name:
            return "RECTANGLE"
        elif "EllipseShape" in service_name:
            return "ELLIPSE"
        elif "ConnectorShape" in service_name:
            return "CONNECTOR"
        elif "CustomShape" in service_name:
            return "CUSTOM_SHAPE"
        elif "ControlShape" in service_name:
            return "CONTROL"
        elif "OLE2Shape" in service_name:
            return "OLE_OBJECT"
        elif "FrameShape" in service_name:
            return "FRAME"
        else:
            return service_name.replace("com.sun.star.drawing.", "").upper()
    except Exception:
        return "UNKNOWN"

def _extract_text_frame_info(shape: Any) -> Dict[str, Any]:
    """
    提取文本框的详细信息，包括段落和字体。
    """
    text_frame_info = {
        "paragraphs": []
    }
    try:
        text_cursor = shape.createTextCursor()
        text_cursor.gotoStart(False)
        
        # 遍历段落
        while True:
            paragraph = text_cursor.getText().getStartOfParagraph(text_cursor)
            if not paragraph:
                break
            
            para_info = {
                "text": paragraph.getString(),
                "runs": []
            }
            
            # 遍历段落中的文本片段 (runs)
            run_cursor = paragraph.getText().createTextCursorByRange(paragraph)
            run_cursor.gotoStart(False)
            
            while True:
                run_cursor.gotoEndOfWord(False) # 移动到单词末尾
                run_text = run_cursor.getString()
                if not run_text:
                    break
                
                run_info = {
                    "text": run_text,
                    "font_name": getattr(run_cursor, 'CharFontName', None),
                    "font_size_pt": getattr(run_cursor, 'CharHeight', None),
                    "bold": getattr(run_cursor, 'CharWeight', FontWeight.NORMAL) == FontWeight.BOLD,
                    "italic": getattr(run_cursor, 'CharPosture', FontSlant.NONE) != FontSlant.NONE,
                    "underline": getattr(run_cursor, 'CharUnderline', 0) != 0,
                    "superscript": getattr(run_cursor, 'CharSuperScript', False),
                    "subscript": getattr(run_cursor, 'CharSubScript', False),
                    "font_color_rgb": getattr(run_cursor, 'CharColor', None),
                    "background_color_rgb": getattr(run_cursor, 'CharBackColor', None),
                }
                para_info["runs"].append(run_info)
                
                if not run_cursor.gotoNextWord(False):
                    break
            
            text_frame_info["paragraphs"].append(para_info)
            
            if not text_cursor.gotoNextParagraph(False):
                break
                
    except Exception as e:
        logger.debug(f"提取文本框信息失败: {e}")
    return text_frame_info

def _extract_fill_info(shape: Any) -> Dict[str, Any]:
    """提取形状的填充信息。"""
    fill_info = {}
    try:
        fill_info["style"] = shape.FillStyle.value
        if shape.FillStyle == FillStyle.SOLID:
            fill_info["color_rgb"] = shape.FillColor
        elif shape.FillStyle == FillStyle.GRADIENT:
            fill_info["gradient_start_color_rgb"] = shape.FillGradientColor
            fill_info["gradient_end_color_rgb"] = shape.FillGradientEndColor
            # 更多渐变属性...
        elif shape.FillStyle == FillStyle.BITMAP:
            fill_info["bitmap_url"] = shape.FillBitmapURL
            # 更多位图属性...
    except Exception as e:
        logger.debug(f"提取填充信息失败: {e}")
    return fill_info

def _extract_line_info(shape: Any) -> Dict[str, Any]:
    """提取形状的线条信息。"""
    line_info = {}
    try:
        line_info["style"] = shape.LineStyle.value
        if shape.LineStyle != LineStyle.NONE:
            line_info["color_rgb"] = shape.LineColor
            line_info["width_pt"] = shape.LineWidth / 100 # 转换为磅 (1/100 mm)
            line_info["transparency"] = shape.LineTransparence
            # 更多线条属性...
    except Exception as e:
        logger.debug(f"提取线条信息失败: {e}")
    return line_info

if __name__ == '__main__':
    # 这是一个简单的测试用例，需要一个实际的PPT文件路径
    # 请替换为您的PPT文件路径
    test_ppt_file = "D:/project/FCI/test_presentation.pptx" # 假设存在这个文件

    # 创建一个简单的PPT文件用于测试 (如果不存在)
    if not os.path.exists(test_ppt_file):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6] # Blank layout
            
            # Slide 1
            slide1 = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1.0)
            txBox = slide1.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "Hello, UNO!"
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(24)
            p.font.name = "Arial"
            p.font.color.rgb = (0xFF, 0x00, 0x00) # Red
            
            left = Inches(3.0)
            shape1 = slide1.shapes.add_shape(MSO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            shape1.text = "Rectangle Shape"
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = (0x00, 0x00, 0xFF) # Blue
            shape1.line.color.rgb = (0x00, 0xFF, 0x00) # Green
            shape1.line.width = Pt(2)

            # Slide 2
            slide2 = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1.0)
            txBox2 = slide2.shapes.add_textbox(left, top, width, height)
            tf2 = txBox2.text_frame
            tf2.text = "Second Slide Text"
            
            prs.save(test_ppt_file)
            print(f"Created a test PPT file: {test_ppt_file}")
        except ImportError:
            print("python-pptx not installed. Cannot create test PPT. Please provide an existing PPT file.")
        except Exception as e:
            print(f"Error creating test PPT: {e}")

    if os.path.exists(test_ppt_file):
        print(f"尝试获取PPT格式信息: {test_ppt_file}")
        
        # 获取所有幻灯片的格式信息
        all_format_data = get_ppt_format_uno(test_ppt_file)
        import json
        print("\n--- All Slides Format Info ---")
        print(json.dumps(all_format_data, indent=2, ensure_ascii=False))

        # 获取第一页的格式信息
        first_slide_format_data = get_ppt_format_uno(test_ppt_file, slide_index=0)
        print("\n--- First Slide Format Info ---")
        print(json.dumps(first_slide_format_data, indent=2, ensure_ascii=False))
    else:
        print("未找到测试PPT文件，请手动创建或指定一个有效的PPT文件路径。")
